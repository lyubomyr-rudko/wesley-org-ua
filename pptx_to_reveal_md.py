#!/usr/bin/env python3
import argparse
import re
from pathlib import Path
from zipfile import ZipFile
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def clean(text: str) -> str:
    text = text.replace("\x0b", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def escape_md(text: str) -> str:
    return text.replace("|", "\\|")


def shape_text(shape):
    chunks = []

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            chunks.extend(shape_text(s))
        return chunks

    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
        txt = clean(shape.text)
        if txt:
            chunks.append(txt)

    if hasattr(shape, "has_table") and shape.has_table:
        table = shape.table
        rows = []
        for row in table.rows:
            rows.append([escape_md(clean(cell.text)) for cell in row.cells])

        if rows:
            header = "| " + " | ".join(rows[0]) + " |"
            sep = "| " + " | ".join(["---"] * len(rows[0])) + " |"
            body = ["| " + " | ".join(r) + " |" for r in rows[1:]]
            chunks.append("\n".join([header, sep] + body))

    return chunks


def extract_notes(pptx_path: Path, slide_num: int) -> str:
    notes_path = f"ppt/notesSlides/notesSlide{slide_num}.xml"

    try:
        with ZipFile(pptx_path) as z:
            if notes_path not in z.namelist():
                return ""
            xml = z.read(notes_path)
    except Exception:
        return ""

    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

    try:
        root = ET.fromstring(xml)
        texts = [t.text for t in root.findall(".//a:t", ns) if t.text]
        return clean("\n".join(texts))
    except Exception:
        return ""


def convert(pptx_path: Path, out_path: Path, include_notes: bool):
    prs = Presentation(pptx_path)
    slides_md = []

    for idx, slide in enumerate(prs.slides, start=1):
        shapes = sorted(
            slide.shapes,
            key=lambda s: (getattr(s, "top", 0), getattr(s, "left", 0)),
        )

        chunks = []
        for shape in shapes:
            chunks.extend(shape_text(shape))

        slide_md = "\n\n".join(chunks).strip()

        if include_notes:
            notes = extract_notes(pptx_path, idx)
            if notes:
                slide_md += f"\n\nNotes:\n{notes}"

        slides_md.append(slide_md)

    out_path.write_text("\n\n---\n\n".join(slides_md), encoding="utf-8")


def main():
    parser = argparse.ArgumentParser(
        description="Convert PPTX text to reveal.js-compatible Markdown."
    )
    parser.add_argument("input", help="Input .pptx file")
    parser.add_argument("-o", "--output", help="Output .md file")
    parser.add_argument("--no-notes", action="store_true", help="Do not include speaker notes")

    args = parser.parse_args()

    pptx_path = Path(args.input)
    out_path = Path(args.output) if args.output else pptx_path.with_suffix(".md")

    convert(pptx_path, out_path, include_notes=not args.no_notes)
    print(f"Wrote {out_path}")


if __name__ == "__main__":
    main()